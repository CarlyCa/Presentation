from flask import Flask, render_template, jsonify, request, send_file
import json
import os
from pptx import Presentation
import requests
from flask_cors import CORS

# Flask app
app = Flask(__name__)

# Enable CORS for all routes and origins
CORS(app, resources={r"/*": {"origins": "*"}}, 
     supports_credentials=True,
     allow_headers=["Content-Type", "Authorization"],
     methods=["GET", "POST", "OPTIONS"])


# Constants
API_URL = "https://api.stack-ai.com/inference/v0/run/ac522f2a-ccb2-4608-8bab-3e1ccf74af42/6744a3cc842d4ea48dda5ddf"
HEADERS = {
    'Authorization': 'Bearer f39536ea-6cd4-4afe-ac99-bd9babd2cd21',
    'Content-Type': 'application/json'
}
OUTPUT_DIR = './output'
TEMPLATE_PATH = 'template.pptx'

# Ensure output directory exists
os.makedirs(OUTPUT_DIR, exist_ok=True)

@app.route("/")
def home():
    """Serve the HTML page."""
    return render_template("index.html")

def query_api(payload):
    """Fetch slide content from API."""
    response = requests.post(API_URL, headers=HEADERS, json=payload)
    if response.status_code != 200:
        return {"error": f"API request failed with status code {response.status_code}"}

    try:
        data = response.json()
        raw_out1 = data.get("outputs", {}).get("out-1", None)
        if raw_out1:
            if raw_out1.startswith("```json") and raw_out1.endswith("```"):
                raw_out1 = raw_out1[7:-3]
            slides = json.loads(raw_out1)
            return slides.get("slides", [])
        else:
            return {"error": "'out-1' key not found in the API response."}
    except (json.JSONDecodeError, KeyError) as e:
        return {"error": f"Error parsing API response: {e}"}

def create_slides_from_content(slide_content):
    """Create a PowerPoint presentation from slide content."""
    if not os.path.exists(TEMPLATE_PATH):
        return {"error": f"Template file '{TEMPLATE_PATH}' does not exist."}

    prs = Presentation(TEMPLATE_PATH)
    for slide_data in slide_content:
        layout_name = slide_data.get('layout', 'title')
        placeholders = slide_data.get('content', {})
        layout_index = {
            "title": 0,
            "subtitle": 1,
            "3_STEP": 2,
            "4_NUMBERS": 3,
            "5_BULLET_IMAGE": 4,
            "TIMELINE": 5,
            "6_STEPS": 6,
            "Agenda": 7,
            "One_Sentence": 8,
            "3_BULLETS_IMAGE": 9,
            "4_TOPIC": 10
        }.get(layout_name, 0)

        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)

        for placeholder_idx, text in placeholders.items():
            try:
                placeholder_num = int(placeholder_idx.split()[1])
                slide.placeholders[placeholder_num].text = text
            except (ValueError, IndexError, KeyError):
                continue

    output_path = os.path.join(OUTPUT_DIR, 'generated_presentation.pptx')
    prs.save(output_path)
    return {"success": True, "file_path": output_path}

@app.route('/generate', methods=['POST'])
def generate_pptx():   
    if request.method == 'OPTIONS':
        response = make_response()
        response.headers["Access-Control-Allow-Origin"] = "*"
        response.headers["Access-Control-Allow-Methods"] = "POST, OPTIONS"
        response.headers["Access-Control-Allow-Headers"] = "Content-Type"
        return response
    try:
        data = request.json
        if not data or "text" not in data:
            return jsonify({"error": "Missing 'text' field in request payload"}), 400

        slide_content = query_api({"user_id": "12345", "in-0": data["text"]})
        if isinstance(slide_content, dict) and "error" in slide_content:
            return jsonify(slide_content), 500

        result = create_slides_from_content(slide_content)
        if "error" in result:
            return jsonify(result), 500

        return send_file(result["file_path"], as_attachment=True, download_name="presentation.pptx")

    except Exception as e:
        import traceback
        traceback.print_exc()  # Log error details
        return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=int(os.environ.get("PORT", 5000)))
